"""Models to manage graphical attributes of presentation layouts, paragraphs, and charts."""

# Third-Party Libraries
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.util import Pt

# Font Attributes
SMALL = Pt(10)
LARGE = Pt(28)
BLUE = RGBColor(3, 37, 126)
GREEN = RGBColor(20, 200, 100)
WHITE = RGBColor(255, 255, 255)


class Paragraph:
    """Simple class to call text frame attributes."""

    @staticmethod
    def shapes(slide):
        """Create a text frame."""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
        return shape

    @staticmethod
    def shapes_find(slide):
        """Find text frames."""
        for shape in slide.shapes:
            print(shape.name)
        return shape

    @staticmethod
    def text_frame(shape):
        """Clear the content of a text frame."""
        text_frame = shape.text_frame
        text_frame.clear()
        return text_frame

    @staticmethod
    def text_frame_ov_val(slide, name):
        """Text style for overview page values."""
        for shape in slide.shapes:
            if shape.name == name:
                text_frame = shape.text_frame
                text_frame.clear()
        return text_frame

    @staticmethod
    def text_frame_key_metric(slide, name):
        """Text style for metrics."""
        for shape in slide.shapes:
            if shape.name == name:
                text_frame = shape.text_frame
                text_frame.clear()
        return text_frame

    @staticmethod
    def text_style_title(font):
        """Text style for cover page title."""
        font.name = "Calibri"
        font.size, font.color = LARGE, WHITE
        return font

    @staticmethod
    def text_style_key_metric(font):
        """Text style for key metrics."""
        font.name = "Calibri"
        font.size, font.color = SMALL, BLUE
        return font

    @staticmethod
    def text_style_ov_val(font):
        """Text style for overview page values."""
        font.name = "Calibri"
        font.size, font.color = LARGE, BLUE
        return font


class Graph:
    """Simple class to call chart attributes."""

    @staticmethod
    def chart_med(chart):
        """Medium chart."""
        chart.font.size, chart.font.color = SMALL, GREEN
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    @staticmethod
    def chart_sm(chart):
        """Small chart."""
        chart.font.size, chart.font.color = SMALL, GREEN
